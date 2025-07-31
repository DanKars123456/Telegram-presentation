import os
import logging
import telebot
from telebot import types
from pptx import Presentation
from pptx.util import Inches
from openai import OpenAI
from openai.types.chat import ChatCompletionMessage
from openai.types.chat import ChatCompletion
from openai import APIConnectionError, RateLimitError, BadRequestError

# Логирование
logging.basicConfig(level=logging.INFO)

# Токен Telegram
TOKEN = os.getenv("8237141546:AAH4Duh0H3B7kjGn5hfu5Z4DbiaxjBY383c") or "your_telegram_token_here"
bot = telebot.TeleBot(TOKEN)

# Токен OpenAI
client = OpenAI(api_key=os.getenv("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") or "your_openai_api_key_here")

# Состояние пользователей
user_state = {}

# Шаги взаимодействия
steps = ["name", "topic", "slides", "design", "font", "style", "language"]

# Команда /start
@bot.message_handler(commands=["start"])
def handle_start(message):
    user_id = message.chat.id
    user_state[user_id] = {"step": "name"}
    bot.send_message(user_id, "👋 Привет! Как тебя зовут?")

# Ответы на текст
@bot.message_handler(func=lambda msg: True)
def handle_message(message):
    user_id = message.chat.id
    text = message.text
    state = user_state.get(user_id, {"step": "name"})

    step = state["step"]

    if step == "name":
        state["name"] = text
        state["step"] = "topic"
        bot.send_message(user_id, f"Приятно познакомиться, {text}! Напиши тему презентации:")
    elif step == "topic":
        state["topic"] = text
        state["step"] = "slides"
        bot.send_message(user_id, "Сколько слайдов ты хочешь?")
    elif step == "slides":
        if not text.isdigit():
            bot.send_message(user_id, "Введите число.")
            return
        state["slides"] = int(text)
        state["step"] = "design"
        markup = types.ReplyKeyboardMarkup(resize_keyboard=True)
        markup.add("Светлый минимализм", "Темный контрастный")
        bot.send_message(user_id, "Выбери стиль оформления:", reply_markup=markup)
    elif step == "design":
        state["design"] = text
        state["step"] = "font"
        markup = types.ReplyKeyboardMarkup(resize_keyboard=True)
        markup.add("Arial", "Times New Roman")
        bot.send_message(user_id, "Выбери шрифт:", reply_markup=markup)
    elif step == "font":
        state["font"] = text
        state["step"] = "style"
        markup = types.ReplyKeyboardMarkup(resize_keyboard=True)
        markup.add("Только текст", "Текст с изображением")
        bot.send_message(user_id, "Выбери тип слайдов:", reply_markup=markup)
    elif step == "style":
        state["style"] = text
        state["step"] = "language"
        markup = types.ReplyKeyboardMarkup(resize_keyboard=True)
        markup.add("Русский", "English", "Polski")
        bot.send_message(user_id, "Выбери язык презентации:", reply_markup=markup)
    elif step == "language":
        state["language"] = text
        bot.send_message(user_id, "⏳ Генерирую презентацию...")
        try:
            pptx_file = generate_presentation(state)
            with open(pptx_file, "rb") as f:
                bot.send_document(user_id, f)
            os.remove(pptx_file)
        except Exception as e:
            bot.send_message(user_id, f"❌ Ошибка при создании презентации: {e}")
        finally:
            user_state.pop(user_id, None)
    user_state[user_id] = state


def generate_presentation(state):
    topic = state["topic"]
    slides = state["slides"]
    language = state["language"]

    prompt = f"Create a presentation on the topic '{topic}' in {language}. Include {slides} slides. Use bullet points."

    try:
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{
                "role": "user",
                "content": prompt
            }],
            temperature=0.7,
        )
        content = response.choices[0].message.content
    except (APIConnectionError, RateLimitError, BadRequestError) as e:
        raise RuntimeError(f"OpenAI error: {e}")

    prs = Presentation()
    for i, slide_text in enumerate(content.split("\n\n")):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title, *points = slide_text.split("\n")
        slide.shapes.title.text = title.strip()
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        for point in points:
            tf.add_paragraph().text = point.strip("• ").strip()

    filename = f"presentation_{state['name']}.pptx"
    prs.save(filename)
    return filename


if __name__ == "__main__":
    try:
        print("✅ Бот запускается...")
        bot.infinity_polling()
    except Exception as e:
        print(f"❌ Бот не запустился: {e}")
